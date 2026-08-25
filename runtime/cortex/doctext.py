"""Extract readable text from office documents so the drafter can READ what clients attach.

Images and PDFs go to the model natively (provider.media_blocks); everything else comes through
here as extracted text. Pure-python extractors (python-docx / openpyxl / xlrd / python-pptx) —
no LibreOffice, nothing shells out. A format we can't read returns '' and the caller says so
honestly on the card rather than pretending.
"""
from __future__ import annotations

import csv
import io

MAX_CHARS = 15_000          # per file — plenty to reply to a brief/quote; keeps the prompt sane
_MAX_ROWS_PER_SHEET = 200   # spreadsheets: enough for any client-facing table, not a data dump

# canonical office mimes + the extensions used when Gmail sends a generic octet-stream
DOC_MIMES = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "application/vnd.ms-excel.sheet.macroenabled.12": "xlsx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/vnd.ms-excel": "xls",
    "text/csv": "csv",
    "text/plain": "txt",
}
DOC_EXTS = {"docx", "xlsx", "xlsm", "pptx", "xls", "csv", "txt"}


def kind_for(mime: str, filename: str) -> str | None:
    """Which extractor handles this attachment — by mime first, filename extension as fallback."""
    k = DOC_MIMES.get((mime or "").lower().split(";")[0])
    if k:
        return k
    ext = (filename or "").rsplit(".", 1)[-1].lower()
    return {"xlsm": "xlsx"}.get(ext, ext) if ext in DOC_EXTS else None


def extract(mime: str, filename: str, data: bytes) -> str:
    """Extracted text of the document, trimmed to MAX_CHARS. '' when unreadable — never raises."""
    kind = kind_for(mime, filename)
    try:
        if kind == "docx":
            return _docx(data)[:MAX_CHARS]
        if kind == "xlsx":
            return _xlsx(data)[:MAX_CHARS]
        if kind == "xls":
            return _xls(data)[:MAX_CHARS]
        if kind == "pptx":
            return _pptx(data)[:MAX_CHARS]
        if kind in ("csv", "txt"):
            return _text(data)[:MAX_CHARS]
    except Exception:  # noqa: BLE001 — a corrupt file must never break drafting
        return ""
    return ""


def _text(data: bytes) -> str:
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return ""


def _docx(data: bytes) -> str:
    import docx
    d = docx.Document(io.BytesIO(data))
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    for t in d.tables:
        for row in t.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _xlsx(data: bytes) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"=== Sheet: {ws.title} ===")
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= _MAX_ROWS_PER_SHEET:
                parts.append(f"(… more rows in '{ws.title}' not shown)")
                break
            cells = ["" if v is None else str(v) for v in row]
            if any(c.strip() for c in cells):
                parts.append(" | ".join(cells).rstrip(" |"))
    wb.close()
    return "\n".join(parts)


def _xls(data: bytes) -> str:
    import xlrd
    wb = xlrd.open_workbook(file_contents=data)
    parts = []
    for ws in wb.sheets():
        parts.append(f"=== Sheet: {ws.name} ===")
        for i in range(min(ws.nrows, _MAX_ROWS_PER_SHEET)):
            cells = [str(c.value) if c.value not in ("", None) else "" for c in ws.row(i)]
            if any(c.strip() for c in cells):
                parts.append(" | ".join(cells).rstrip(" |"))
        if ws.nrows > _MAX_ROWS_PER_SHEET:
            parts.append(f"(… more rows in '{ws.name}' not shown)")
    return "\n".join(parts)


def _pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [sh.text_frame.text.strip() for sh in slide.shapes
                 if sh.has_text_frame and sh.text_frame.text.strip()]
        if texts:
            parts.append(f"=== Slide {i} ===\n" + "\n".join(texts))
    return "\n\n".join(parts)


def _csv_preview(data: bytes) -> str:   # kept for completeness; csv routes through _text today
    rows = list(csv.reader(io.StringIO(_text(data))))[:_MAX_ROWS_PER_SHEET]
    return "\n".join(" | ".join(r) for r in rows)
